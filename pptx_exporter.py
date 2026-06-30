"""Export recognized diagram elements directly to a single-slide PowerPoint deck.

The exporter consumes the pipeline's element records and OCR text blocks directly.
Editable shapes, text boxes, and straight connectors are emitted as native
PowerPoint objects, while raster icon/picture crops are embedded as images.
"""

import base64
import importlib.util
import io
from pathlib import Path
from typing import Any, Dict, List, Optional, Sequence, Tuple


EMU_PER_PX = 9525  # PowerPoint uses EMUs; 96dpi pixel -> 914400 / 96.


CONNECTOR_TYPES = {"arrow", "line", "connector"}
IMAGE_TYPES = {"icon", "picture", "logo", "chart", "function_graph"}


def is_pptx_export_available() -> bool:
    """Return whether the optional python-pptx runtime dependency is installed."""
    return importlib.util.find_spec("pptx") is not None


def missing_pptx_dependency_message() -> str:
    """Human-friendly install hint for PPTX export."""
    return "PPTX export requires python-pptx. Install it with: pip install python-pptx>=1.0.2"


def export_elements_to_pptx(
    elements: Sequence[Any],
    text_blocks: Optional[Sequence[Dict[str, Any]]],
    canvas_width: int,
    canvas_height: int,
    pptx_path: str,
) -> str:
    """Create a PPTX directly from recognized elements and OCR text blocks."""
    if not is_pptx_export_available():
        raise RuntimeError(missing_pptx_dependency_message())

    from pptx import Presentation
    from pptx.dml.color import RGBColor
    from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE, MSO_CONNECTOR
    from pptx.enum.text import MSO_VERTICAL_ANCHOR, PP_ALIGN
    from pptx.oxml.xmlchemy import OxmlElement
    from pptx.util import Emu, Pt

    prs = Presentation()
    prs.slide_width = Emu(float(canvas_width or 1280) * EMU_PER_PX)
    prs.slide_height = Emu(float(canvas_height or 720) * EMU_PER_PX)
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    for elem in _sort_elements(elements):
        elem_type = str(getattr(elem, "element_type", "") or "").lower()
        if elem_type in CONNECTOR_TYPES:
            _add_connector_from_element(slide, elem, MSO_CONNECTOR, RGBColor, Pt, OxmlElement)
        elif elem_type in IMAGE_TYPES:
            if getattr(elem, "base64", None):
                _add_image_from_element(slide, elem, Emu)
        else:
            _add_shape_from_element(slide, elem, MSO_AUTO_SHAPE_TYPE, Emu, Pt, RGBColor)

    for block in text_blocks or []:
        _add_text_block(slide, block, Emu, Pt, RGBColor, MSO_VERTICAL_ANCHOR, PP_ALIGN)

    output_path = Path(pptx_path)
    output_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(output_path)
    return str(output_path)


def _sort_elements(elements: Sequence[Any]) -> List[Any]:
    def sort_key(elem: Any) -> Tuple[int, int]:
        layer = int(getattr(elem, "layer_level", 5) or 5)
        bbox = getattr(elem, "bbox", None)
        area = int(getattr(bbox, "area", 0) or 0)
        return layer, -area

    return sorted(elements or [], key=sort_key)


def _emu(value: float, Emu):
    return Emu(float(value) * EMU_PER_PX)


def _rgb(color: Optional[str], RGBColor):
    if not color or color == "none" or not str(color).startswith("#") or len(str(color)) != 7:
        return None
    text = str(color)
    return RGBColor(int(text[1:3], 16), int(text[3:5], 16), int(text[5:7], 16))


def _shape_type(elem_type: str, MSO_AUTO_SHAPE_TYPE):
    normalized = elem_type.replace("_", " ")
    if normalized == "rounded rectangle":
        return MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE
    mapping = {
        "ellipse": "OVAL",
        "circle": "OVAL",
        "cylinder": "CAN",
        "database": "CAN",
        "diamond": "DIAMOND",
        "triangle": "TRIANGLE",
        "hexagon": "HEXAGON",
        "parallelogram": "PARALLELOGRAM",
        "cloud": "CLOUD",
    }
    shape_name = mapping.get(normalized, "RECTANGLE")
    return getattr(MSO_AUTO_SHAPE_TYPE, shape_name, MSO_AUTO_SHAPE_TYPE.RECTANGLE)


def _add_shape_from_element(slide, elem: Any, MSO_AUTO_SHAPE_TYPE, Emu, Pt, RGBColor) -> None:
    bbox = getattr(elem, "bbox", None)
    if bbox is None or getattr(bbox, "width", 0) <= 0 or getattr(bbox, "height", 0) <= 0:
        return
    elem_type = str(getattr(elem, "element_type", "rectangle") or "rectangle").lower()
    shape = slide.shapes.add_shape(
        _shape_type(elem_type, MSO_AUTO_SHAPE_TYPE),
        _emu(getattr(bbox, "x1", 0), Emu),
        _emu(getattr(bbox, "y1", 0), Emu),
        _emu(getattr(bbox, "width", 0), Emu),
        _emu(getattr(bbox, "height", 0), Emu),
    )

    fill = _rgb(getattr(elem, "fill_color", None) or "#ffffff", RGBColor)
    if fill is None:
        shape.fill.background()
    else:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill

    stroke = _rgb(getattr(elem, "stroke_color", None) or "#000000", RGBColor)
    if stroke is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = stroke
        shape.line.width = Pt(float(getattr(elem, "stroke_width", 1) or 1))


def _add_image_from_element(slide, elem: Any, Emu) -> None:
    bbox = getattr(elem, "bbox", None)
    image_data = getattr(elem, "base64", None)
    if bbox is None or not image_data:
        return
    image_bytes = base64.b64decode(image_data)
    slide.shapes.add_picture(
        io.BytesIO(image_bytes),
        _emu(getattr(bbox, "x1", 0), Emu),
        _emu(getattr(bbox, "y1", 0), Emu),
        _emu(getattr(bbox, "width", 0), Emu),
        _emu(getattr(bbox, "height", 0), Emu),
    )


def _add_connector_from_element(slide, elem: Any, MSO_CONNECTOR, RGBColor, Pt, OxmlElement) -> None:
    start, end = _connector_points(elem)
    if start is None or end is None:
        return
    connector = slide.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT,
        int(float(start[0]) * EMU_PER_PX),
        int(float(start[1]) * EMU_PER_PX),
        int(float(end[0]) * EMU_PER_PX),
        int(float(end[1]) * EMU_PER_PX),
    )
    stroke = _rgb(getattr(elem, "stroke_color", None) or "#000000", RGBColor)
    if stroke is not None:
        connector.line.color.rgb = stroke
    connector.line.width = Pt(float(getattr(elem, "stroke_width", 2) or 2))
    if getattr(elem, "line_style", None) in {"dashed", "dotted"} or getattr(elem, "dash_pattern", None):
        _apply_connector_dash(connector, OxmlElement)
    _apply_connector_arrowheads(connector, getattr(elem, "arrow_heads", None), getattr(elem, "element_type", ""), OxmlElement)


def _connector_points(elem: Any) -> Tuple[Optional[Tuple[float, float]], Optional[Tuple[float, float]]]:
    start = getattr(elem, "arrow_start", None)
    end = getattr(elem, "arrow_end", None)
    if start and end:
        return start, end

    bbox = getattr(elem, "bbox", None)
    if bbox is None:
        return None, None
    cx = (getattr(bbox, "x1", 0) + getattr(bbox, "x2", 0)) / 2
    cy = (getattr(bbox, "y1", 0) + getattr(bbox, "y2", 0)) / 2
    if getattr(bbox, "width", 0) >= getattr(bbox, "height", 0):
        return (float(getattr(bbox, "x1", 0)), float(cy)), (float(getattr(bbox, "x2", 0)), float(cy))
    return (float(cx), float(getattr(bbox, "y1", 0))), (float(cx), float(getattr(bbox, "y2", 0)))


def _add_text_block(slide, block: Dict[str, Any], Emu, Pt, RGBColor, MSO_VERTICAL_ANCHOR, PP_ALIGN) -> None:
    geo = block.get("geometry") or {}
    x = float(geo.get("x", 0))
    y = float(geo.get("y", 0))
    width = max(float(geo.get("width", 20)), 20.0)
    height = max(float(geo.get("height", 10)), 10.0)
    textbox = slide.shapes.add_textbox(_emu(x, Emu), _emu(y, Emu), _emu(width, Emu), _emu(height, Emu))
    frame = textbox.text_frame
    frame.clear()
    frame.margin_left = 0
    frame.margin_right = 0
    frame.margin_top = 0
    frame.margin_bottom = 0
    frame.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE
    frame.word_wrap = True
    paragraph = frame.paragraphs[0]
    paragraph.alignment = PP_ALIGN.CENTER
    run = paragraph.add_run()
    run.text = str(block.get("text", ""))
    run.font.size = Pt(float(block.get("font_size", 12) or 12))
    run.font.bold = bool(block.get("is_bold") or str(block.get("font_weight", "")).lower() == "bold")
    run.font.italic = bool(block.get("is_italic") or str(block.get("font_style", "")).lower() == "italic")
    font_color = _rgb(block.get("font_color") or "#000000", RGBColor)
    if font_color is not None:
        run.font.color.rgb = font_color
    if block.get("font_family"):
        run.font.name = str(block["font_family"])


def _apply_connector_dash(connector, OxmlElement) -> None:
    line = connector._element.spPr.get_or_add_ln()
    for child in list(line):
        if child.tag.endswith("}prstDash"):
            line.remove(child)
    dash = OxmlElement("a:prstDash")
    dash.set("val", "dash")
    line.append(dash)


def _apply_connector_arrowheads(connector, arrow_heads: Optional[str], elem_type: str, OxmlElement) -> None:
    heads = str(arrow_heads or ("end" if str(elem_type).lower() == "arrow" else "none")).lower()
    if heads in {"none", "null", ""}:
        return
    line = connector._element.spPr.get_or_add_ln()
    if heads in {"start", "both"}:
        head = OxmlElement("a:headEnd")
        head.set("type", "triangle")
        line.append(head)
    if heads in {"end", "both"}:
        tail = OxmlElement("a:tailEnd")
        tail.set("type", "triangle")
        line.append(tail)
